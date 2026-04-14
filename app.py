import io
import streamlit as st
from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData

st.set_page_config(page_title="Weekly Report Generator", layout="centered")
st.title("Weekly Report Generator")

excel_file = st.file_uploader("ارفعي ملف Excel", type=["xlsx"])
ppt_file = st.file_uploader("ارفعي الباوربوينت الريفرنس", type=["pptx"])


# =========================
# General Helpers
# =========================
def get_sheet_case_insensitive(wb, target_name):
    for sheet_name in wb.sheetnames:
        if sheet_name.strip().lower() == target_name.strip().lower():
            return wb[sheet_name]
    return None


def get_chart_shapes(slide):
    return [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]


def trailing_zeros_to_none(values):
    result = list(values)

    last_non_zero_index = -1
    for i, v in enumerate(result):
        if v not in (None, 0, 0.0):
            last_non_zero_index = i

    if last_non_zero_index == -1:
        return [None for _ in result]

    for i in range(last_non_zero_index + 1, len(result)):
        result[i] = None

    return result


def normalize_percent(v):
    if v is None or v == "":
        return 0.0
    return float(v) / 100.0


def filter_rows_by_type(ws, row_type_name):
    rows = []
    for row in range(2, ws.max_row + 1):
        row_type = ws.cell(row=row, column=1).value
        if str(row_type).strip().lower() == row_type_name.strip().lower():
            rows.append(row)
    return rows


def sort_three_charts_layout(charts):
    # top chart first, then bottom-left, then bottom-right
    return sorted(charts, key=lambda s: (s.top, s.left))


def set_percent_axis_and_labels(chart, decimals=2):
    fmt = "0." + ("0" * decimals) + "%"

    try:
        chart.value_axis.tick_labels.number_format = fmt
    except Exception:
        pass

    for series in chart.series:
        try:
            series.data_labels.show_value = True
            series.data_labels.number_format = fmt
        except Exception:
            pass


def set_plain_number_axis_and_labels(chart):
    try:
        chart.value_axis.tick_labels.number_format = '#,##0'
    except Exception:
        pass

    for series in chart.series:
        try:
            series.data_labels.show_value = False
        except Exception:
            pass


def replace_single_series_chart(chart, categories, series_name, values, percent_chart=False, percent_decimals=2):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    chart.replace_data(chart_data)

    if percent_chart:
        set_percent_axis_and_labels(chart, decimals=percent_decimals)
    else:
        set_plain_number_axis_and_labels(chart)


def replace_two_series_chart(chart, categories, series1_name, values1, series2_name, values2, percent_chart=False, percent_decimals=2):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series1_name, values1)
    chart_data.add_series(series2_name, values2)
    chart.replace_data(chart_data)

    if percent_chart:
        set_percent_axis_and_labels(chart, decimals=percent_decimals)
    else:
        set_plain_number_axis_and_labels(chart)


def read_values_from_rows(ws, rows, columns, percent_keys=None):
    percent_keys = percent_keys or set()
    result = {key: [] for key in columns}

    for row in rows:
        for key, col in columns.items():
            value = ws.cell(row=row, column=col).value

            if key == "week":
                result[key].append(str(value) if value else "")
            elif key in percent_keys:
                result[key].append(normalize_percent(value))
            else:
                result[key].append(float(value) if value is not None else 0)

    return result


# =========================
# STRIP
# Slides 3-4
# Sheet: Strip
# A Week
# B Production Roll
# C Achieved %
# D Slag %
# E Slag kg
# F Target Slag %
# =========================
def read_strip_data(excel_bytes):
    wb = load_workbook(io.BytesIO(excel_bytes), data_only=True)
    ws = get_sheet_case_insensitive(wb, "Strip")

    if ws is None:
        raise ValueError("لا توجد شيت باسم Strip في ملف Excel.")

    weeks = []
    production = []
    achieved = []
    slag_pct = []
    slag_kg = []
    target_slag_pct = []

    for row in range(2, 7):  # Week 1 to Week 5
        week = ws.cell(row=row, column=1).value
        prod = ws.cell(row=row, column=2).value
        ach = ws.cell(row=row, column=3).value
        slg_pct = ws.cell(row=row, column=4).value
        slg_kg_val = ws.cell(row=row, column=5).value
        target_val = ws.cell(row=row, column=6).value

        weeks.append(str(week) if week else f"Week {row-1}")
        production.append(float(prod) if prod is not None else 0)
        achieved.append(normalize_percent(ach))
        slag_pct.append(normalize_percent(slg_pct))
        slag_kg.append(float(slg_kg_val) if slg_kg_val is not None else 0)
        target_slag_pct.append(normalize_percent(target_val) if target_val is not None else 0.028)

    return weeks, production, achieved, slag_pct, slag_kg, target_slag_pct


def update_table_first_col_values(table, values):
    max_rows_to_fill = min(len(values), len(table.rows) - 1)
    for i in range(max_rows_to_fill):
        row_idx = i + 1
        val = int(values[i]) if float(values[i]).is_integer() else values[i]
        table.cell(row_idx, 0).text = str(val)


def update_strip_slides(prs, weeks, production, achieved, slag_pct, slag_kg, target_slag_pct):
    slide3 = prs.slides[2]
    slide4 = prs.slides[3]

    slide3_charts = get_chart_shapes(slide3)
    slide4_charts = get_chart_shapes(slide4)
    slide4_tables = [shape for shape in slide4.shapes if getattr(shape, "has_table", False)]

    if len(slide3_charts) < 2:
        raise ValueError("لم أجد شارتين في Slide 3.")
    if len(slide4_charts) < 1:
        raise ValueError("لم أجد شارتًا في Slide 4.")
    if len(slide4_tables) < 1:
        raise ValueError("لم أجد جدولًا في Slide 4.")

    slide3_charts.sort(key=lambda s: s.left)
    slide4_charts.sort(key=lambda s: s.left)
    slide4_tables.sort(key=lambda s: s.left)

    production_plot = trailing_zeros_to_none(production)
    achieved_plot = trailing_zeros_to_none(achieved)
    slag_pct_plot = trailing_zeros_to_none(slag_pct)

    replace_single_series_chart(
        slide3_charts[0].chart,
        weeks,
        "Production Roll",
        production_plot,
        percent_chart=False
    )

    replace_single_series_chart(
        slide3_charts[1].chart,
        weeks,
        "Achieved %",
        achieved_plot,
        percent_chart=True,
        percent_decimals=1
    )

    replace_two_series_chart(
        slide4_charts[0].chart,
        weeks,
        "Slag %",
        slag_pct_plot,
        "Target Slag %",
        target_slag_pct,
        percent_chart=True,
        percent_decimals=2
    )

    update_table_first_col_values(slide4_tables[0].table, slag_kg)


# =========================
# PASTING
# Slides 5-8
# Sheet: Pasting
# A Week
# B Produced Blocks
# C Achieved %
# D Strip Scrap %
# E Strip Scrap Target %
# F Plate Scrap %
# G Plate Scrap Target %
# H Rejected Plates %
# I Rejected Plates Target %
# =========================
def read_pasting_data(excel_bytes):
    wb = load_workbook(io.BytesIO(excel_bytes), data_only=True)
    ws = get_sheet_case_insensitive(wb, "Pasting")

    if ws is None:
        raise ValueError("لا توجد شيت باسم Pasting في ملف Excel.")

    weeks = []
    produced_blocks = []
    achieved_pct = []
    strip_scrap_pct = []
    strip_scrap_target = []
    plate_scrap_pct = []
    plate_scrap_target = []
    rejected_plates_pct = []
    rejected_plates_target = []

    for row in range(2, 7):  # Week 1 to Week 5
        week = ws.cell(row=row, column=1).value
        produced = ws.cell(row=row, column=2).value
        achieved = ws.cell(row=row, column=3).value
        strip_actual = ws.cell(row=row, column=4).value
        strip_target = ws.cell(row=row, column=5).value
        plate_actual = ws.cell(row=row, column=6).value
        plate_target = ws.cell(row=row, column=7).value
        rejected_actual = ws.cell(row=row, column=8).value
        rejected_target = ws.cell(row=row, column=9).value

        weeks.append(str(week) if week else f"Week {row-1}")
        produced_blocks.append(float(produced) if produced is not None else 0)
        achieved_pct.append(normalize_percent(achieved))
        strip_scrap_pct.append(normalize_percent(strip_actual))
        strip_scrap_target.append(normalize_percent(strip_target) if strip_target is not None else 0.003)
        plate_scrap_pct.append(normalize_percent(plate_actual))
        plate_scrap_target.append(normalize_percent(plate_target) if plate_target is not None else 0.003)
        rejected_plates_pct.append(normalize_percent(rejected_actual))
        rejected_plates_target.append(normalize_percent(rejected_target) if rejected_target is not None else 0.0003)

    return (
        weeks,
        produced_blocks,
        achieved_pct,
        strip_scrap_pct,
        strip_scrap_target,
        plate_scrap_pct,
        plate_scrap_target,
        rejected_plates_pct,
        rejected_plates_target,
    )


def update_pasting_slides(
    prs,
    weeks,
    produced_blocks,
    achieved_pct,
    strip_scrap_pct,
    strip_scrap_target,
    plate_scrap_pct,
    plate_scrap_target,
    rejected_plates_pct,
    rejected_plates_target,
):
    slide5 = prs.slides[4]
    slide6 = prs.slides[5]
    slide7 = prs.slides[6]
    slide8 = prs.slides[7]

    slide5_charts = get_chart_shapes(slide5)
    slide6_charts = get_chart_shapes(slide6)
    slide7_charts = get_chart_shapes(slide7)
    slide8_charts = get_chart_shapes(slide8)

    if len(slide5_charts) < 2:
        raise ValueError("لم أجد شارتين في Slide 5.")
    if len(slide6_charts) < 1:
        raise ValueError("لم أجد شارتًا في Slide 6.")
    if len(slide7_charts) < 1:
        raise ValueError("لم أجد شارتًا في Slide 7.")
    if len(slide8_charts) < 1:
        raise ValueError("لم أجد شارتًا في Slide 8.")

    slide5_charts.sort(key=lambda s: s.left)
    slide6_charts.sort(key=lambda s: s.left)
    slide7_charts.sort(key=lambda s: s.left)
    slide8_charts.sort(key=lambda s: s.left)

    produced_blocks_plot = trailing_zeros_to_none(produced_blocks)
    achieved_pct_plot = trailing_zeros_to_none(achieved_pct)
    strip_scrap_pct_plot = trailing_zeros_to_none(strip_scrap_pct)
    plate_scrap_pct_plot = trailing_zeros_to_none(plate_scrap_pct)
    rejected_plates_pct_plot = trailing_zeros_to_none(rejected_plates_pct)

    replace_single_series_chart(
        slide5_charts[0].chart,
        weeks,
        "Produced Blocks",
        produced_blocks_plot,
        percent_chart=False
    )

    replace_single_series_chart(
        slide5_charts[1].chart,
        weeks,
        "Achieved %",
        achieved_pct_plot,
        percent_chart=True,
        percent_decimals=1
    )

    replace_two_series_chart(
        slide6_charts[0].chart,
        weeks,
        "Strip Scrap %",
        strip_scrap_pct_plot,
        "Target Strip Scrap %",
        strip_scrap_target,
        percent_chart=True,
        percent_decimals=2
    )

    replace_two_series_chart(
        slide7_charts[0].chart,
        weeks,
        "Plate Scrap %",
        plate_scrap_pct_plot,
        "Target Plate Scrap %",
        plate_scrap_target,
        percent_chart=True,
        percent_decimals=2
    )

    replace_two_series_chart(
        slide8_charts[0].chart,
        weeks,
        "Rejected Plates %",
        rejected_plates_pct_plot,
        "Target Rejected Plates %",
        rejected_plates_target,
        percent_chart=True,
        percent_decimals=2
    )


# =========================
# ASSEMBLY MAIN
# Slides 9-12
# Sheet: Assembly_Main
# A Row Type
# B Week
# C Production Battery
# D Productivity %
# =========================
def read_assembly_main_data(excel_bytes):
    wb = load_workbook(io.BytesIO(excel_bytes), data_only=True)
    ws = get_sheet_case_insensitive(wb, "Assembly_Main")

    if ws is None:
        raise ValueError("لا توجد شيت باسم Assembly_Main.")

    column_map = {
        "week": 2,
        "production": 3,
        "productivity": 4,
    }

    percent_keys = {"productivity"}

    return {
        "total": read_values_from_rows(ws, filter_rows_by_type(ws, "Total"), column_map, percent_keys),
        "kory1": read_values_from_rows(ws, filter_rows_by_type(ws, "Kory1"), column_map, percent_keys),
        "kory2": read_values_from_rows(ws, filter_rows_by_type(ws, "Kory2"), column_map, percent_keys),
        "kory3": read_values_from_rows(ws, filter_rows_by_type(ws, "Kory3"), column_map, percent_keys),
    }


def update_assembly_main_slides(prs, main_data):
    slide_map = [
        (prs.slides[8],  main_data["total"]),
        (prs.slides[9],  main_data["kory1"]),
        (prs.slides[10], main_data["kory2"]),
        (prs.slides[11], main_data["kory3"]),
    ]

    for slide, data in slide_map:
        charts = get_chart_shapes(slide)

        if len(charts) < 2:
            raise ValueError("أحد سلايدات 9-12 لا يحتوي على شارتين.")

        charts.sort(key=lambda s: s.left)

        weeks = data["week"]
        production_plot = trailing_zeros_to_none(data["production"])
        productivity_plot = trailing_zeros_to_none(data["productivity"])

        replace_single_series_chart(
            charts[0].chart,
            weeks,
            "Production Battery",
            production_plot,
            percent_chart=False
        )

        replace_single_series_chart(
            charts[1].chart,
            weeks,
            "Productivity %",
            productivity_plot,
            percent_chart=True,
            percent_decimals=1
        )


# =========================
# ASSEMBLY SCRAP
# Slides 13-22
# Sheet: Assembly_Scrap
# A Row Type
# B Week
# C Scraped Plate %
# D Scraped Plate Target %
# E Reworked Plate %
# F Reworked Plate Target %
# G Separator Scrap %
# H Separator Target %
# I Box Scrap %
# J Box Scrap Target %
# K Cover Scrap %
# L Cover Scrap Target %
# =========================
def read_assembly_scrap_data(excel_bytes):
    wb = load_workbook(io.BytesIO(excel_bytes), data_only=True)
    ws = get_sheet_case_insensitive(wb, "Assembly_Scrap")

    if ws is None:
        raise ValueError("لا توجد شيت باسم Assembly_Scrap.")

    column_map = {
        "week": 2,
        "scraped_actual": 3,
        "scraped_target": 4,
        "reworked_actual": 5,
        "reworked_target": 6,
        "separator_actual": 7,
        "separator_target": 8,
        "box_actual": 9,
        "box_target": 10,
        "cover_actual": 11,
        "cover_target": 12,
    }

    percent_keys = {
        "scraped_actual", "scraped_target",
        "reworked_actual", "reworked_target",
        "separator_actual", "separator_target",
        "box_actual", "box_target",
        "cover_actual", "cover_target",
    }

    return {
        "total": read_values_from_rows(ws, filter_rows_by_type(ws, "Total"), column_map, percent_keys),
        "kory1": read_values_from_rows(ws, filter_rows_by_type(ws, "Kory1"), column_map, percent_keys),
        "kory2": read_values_from_rows(ws, filter_rows_by_type(ws, "Kory2"), column_map, percent_keys),
        "kory3": read_values_from_rows(ws, filter_rows_by_type(ws, "Kory3"), column_map, percent_keys),
    }


def update_scrap_total_slide(slide, weeks, actual_vals, target_vals, actual_name, target_name, decimals=2):
    charts = get_chart_shapes(slide)
    if len(charts) < 1:
        raise ValueError("أحد سلايدات التوتال scrap لا يحتوي على شارت.")

    replace_two_series_chart(
        charts[0].chart,
        weeks,
        actual_name,
        trailing_zeros_to_none(actual_vals),
        target_name,
        target_vals,
        percent_chart=True,
        percent_decimals=decimals
    )


def update_scrap_lines_slide(
    slide,
    weeks,
    k1_actual, k1_target,
    k2_actual, k2_target,
    k3_actual, k3_target,
    actual_name, target_name,
    decimals=2
):
    charts = get_chart_shapes(slide)
    if len(charts) < 3:
        raise ValueError("أحد سلايدات line-by-line scrap لا يحتوي على 3 شارتات.")

    charts = sort_three_charts_layout(charts)

    replace_two_series_chart(
        charts[0].chart, weeks,
        actual_name, trailing_zeros_to_none(k1_actual),
        target_name, k1_target,
        percent_chart=True,
        percent_decimals=decimals
    )

    replace_two_series_chart(
        charts[1].chart, weeks,
        actual_name, trailing_zeros_to_none(k2_actual),
        target_name, k2_target,
        percent_chart=True,
        percent_decimals=decimals
    )

    replace_two_series_chart(
        charts[2].chart, weeks,
        actual_name, trailing_zeros_to_none(k3_actual),
        target_name, k3_target,
        percent_chart=True,
        percent_decimals=decimals
    )


def update_assembly_scrap_slides(prs, scrap_data):
    weeks = scrap_data["total"]["week"]

    # Slide 13 / 14 -> Scraped Plate (2 decimals)
    update_scrap_total_slide(
        prs.slides[12], weeks,
        scrap_data["total"]["scraped_actual"],
        scrap_data["total"]["scraped_target"],
        "Scraped Plate %",
        "Target Scraped Plate %",
        decimals=2
    )
    update_scrap_lines_slide(
        prs.slides[13], weeks,
        scrap_data["kory1"]["scraped_actual"], scrap_data["kory1"]["scraped_target"],
        scrap_data["kory2"]["scraped_actual"], scrap_data["kory2"]["scraped_target"],
        scrap_data["kory3"]["scraped_actual"], scrap_data["kory3"]["scraped_target"],
        "Scraped Plate %",
        "Target Scraped Plate %",
        decimals=2
    )

    # Slide 15 / 16 -> Reworked (1 decimal)
    update_scrap_total_slide(
        prs.slides[14], weeks,
        scrap_data["total"]["reworked_actual"],
        scrap_data["total"]["reworked_target"],
        "Reworked Plate %",
        "Target Reworked Plate %",
        decimals=1
    )
    update_scrap_lines_slide(
        prs.slides[15], weeks,
        scrap_data["kory1"]["reworked_actual"], scrap_data["kory1"]["reworked_target"],
        scrap_data["kory2"]["reworked_actual"], scrap_data["kory2"]["reworked_target"],
        scrap_data["kory3"]["reworked_actual"], scrap_data["kory3"]["reworked_target"],
        "Reworked Plate %",
        "Target Reworked Plate %",
        decimals=1
    )

    # Slide 17 / 18 -> Separator (1 decimal)
    update_scrap_total_slide(
        prs.slides[16], weeks,
        scrap_data["total"]["separator_actual"],
        scrap_data["total"]["separator_target"],
        "Separator Scrap %",
        "Target Separator %",
        decimals=1
    )
    update_scrap_lines_slide(
        prs.slides[17], weeks,
        scrap_data["kory1"]["separator_actual"], scrap_data["kory1"]["separator_target"],
        scrap_data["kory2"]["separator_actual"], scrap_data["kory2"]["separator_target"],
        scrap_data["kory3"]["separator_actual"], scrap_data["kory3"]["separator_target"],
        "Separator Scrap %",
        "Target Separator %",
        decimals=1
    )

    # Slide 19 / 20 -> Box (1 decimal)
    update_scrap_total_slide(
        prs.slides[18], weeks,
        scrap_data["total"]["box_actual"],
        scrap_data["total"]["box_target"],
        "Box Scrap %",
        "Target Box Scrap %",
        decimals=1
    )
    update_scrap_lines_slide(
        prs.slides[19], weeks,
        scrap_data["kory1"]["box_actual"], scrap_data["kory1"]["box_target"],
        scrap_data["kory2"]["box_actual"], scrap_data["kory2"]["box_target"],
        scrap_data["kory3"]["box_actual"], scrap_data["kory3"]["box_target"],
        "Box Scrap %",
        "Target Box Scrap %",
        decimals=1
    )

    # Slide 21 / 22 -> Cover (1 decimal)
    update_scrap_total_slide(
        prs.slides[20], weeks,
        scrap_data["total"]["cover_actual"],
        scrap_data["total"]["cover_target"],
        "Cover Scrap %",
        "Target Cover Scrap %",
        decimals=1
    )
    update_scrap_lines_slide(
        prs.slides[21], weeks,
        scrap_data["kory1"]["cover_actual"], scrap_data["kory1"]["cover_target"],
        scrap_data["kory2"]["cover_actual"], scrap_data["kory2"]["cover_target"],
        scrap_data["kory3"]["cover_actual"], scrap_data["kory3"]["cover_target"],
        "Cover Scrap %",
        "Target Cover Scrap %",
        decimals=1
    )


# =========================
# MAIN
# =========================
if st.button("Generate PowerPoint"):
    if excel_file is None or ppt_file is None:
        st.error("ارفعي ملف Excel وملف PowerPoint الأول.")
    else:
        try:
            prs = Presentation(io.BytesIO(ppt_file.getvalue()))

            # Strip
            strip_values = read_strip_data(excel_file.getvalue())
            update_strip_slides(prs, *strip_values)

            # Pasting
            pasting_values = read_pasting_data(excel_file.getvalue())
            update_pasting_slides(prs, *pasting_values)

            # Assembly Main
            assembly_main_values = read_assembly_main_data(excel_file.getvalue())
            update_assembly_main_slides(prs, assembly_main_values)

            # Assembly Scrap
            assembly_scrap_values = read_assembly_scrap_data(excel_file.getvalue())
            update_assembly_scrap_slides(prs, assembly_scrap_values)

            output = io.BytesIO()
            prs.save(output)
            output.seek(0)

            st.success("تم تحديث Slides 3-22 الخاصة بـ Strip وPasting وAssembly بنجاح.")

            st.download_button(
                label="Download PowerPoint",
                data=output,
                file_name="generated_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )

        except Exception as e:
            st.error(f"حصل خطأ: {e}")
